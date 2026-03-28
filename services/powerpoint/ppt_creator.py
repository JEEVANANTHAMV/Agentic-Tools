"""
presentation_creator.py
========================
Hybrid HTML → PowerPoint engine (Screenshot + Editable Text overlay).

Strategy
--------
1. Inject a comprehensive Tailwind CSS utility override so that ANY HTML
   (Tailwind classes, flexbox, grid, arbitrary values, custom CSS) renders
   exactly as it would in a browser — without a running CDN.

2. Use Playwright + Chromium to:
   a. Render the slide at 1280 × 720 px (true browser layout).
   b. Take a full-resolution screenshot → used as the slide background image.
      This captures icons, gradients, shapes, decorative elements PERFECTLY.
   c. Extract every DOM text element's computed position, size, colour,
      font-size, font-weight, font-family, text-align, etc.

3. Reconstruct the slide in python-pptx as:
   - Background image (screenshot)  → covers entire slide — pixel-perfect visuals
   - Text elements                   → transparent textboxes overlaid on top
                                        (selectable / editable in PowerPoint)

4. De-duplicate: skip elements whose text is fully covered by inline children,
   skip Material Icons / font-icon elements, skip purely decorative runs.

Result: the slide looks IDENTICAL to the browser screenshot while keeping
all real text content selectable and editable in PowerPoint.

Dependencies
------------
    pip install playwright python-pptx beautifulsoup4 pillow requests
    playwright install chromium
"""

from __future__ import annotations

import os
import re
import tempfile
import textwrap
from io import BytesIO
from datetime import datetime
from typing import Optional

import requests
from PIL import Image as PILImage
from playwright.async_api import async_playwright

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Canvas ────────────────────────────────────────────────────────────────────
BROWSER_W = 1280   # px  (matches typical ppt-slide class)
BROWSER_H = 720    # px

PPTX_W = Inches(13.333)   # 16:9
PPTX_H = Inches(7.5)

PX_TO_EMU = PPTX_W / BROWSER_W   # EMU per browser-px

# ── Colour helpers ─────────────────────────────────────────────────────────────

def _rgb(r, g, b) -> RGBColor:
    return RGBColor(int(r), int(g), int(b))


def _rgba_to_rgb(lst: Optional[list]) -> Optional[RGBColor]:
    if not lst or len(lst) < 3:
        return None
    return _rgb(*lst[:3])


# ── EMU conversion ─────────────────────────────────────────────────────────────

def _px(v: float) -> Emu:
    """Browser px → python-pptx EMU."""
    return int(v * PX_TO_EMU)


# ── XML helpers ────────────────────────────────────────────────────────────────

def _set_no_line(shape):
    try:
        shape.line.fill.background()
    except Exception:
        pass


# ── Tailwind CSS injector ──────────────────────────────────────────────────────

def _build_tailwind_override_css() -> str:
    """
    Return a <style> block containing CSS equivalents for every common
    Tailwind utility class so that HTML files work correctly when loaded
    from a local file:// URL (where the CDN JIT cannot process classes).
    """
    rules: list[str] = []

    # ── Spacing scale (0-96, step 1 = 4px) ──────────────────────────────────
    for n in range(0, 97):
        v = n * 4
        for prop, cls in [
            ("padding",        f"p-{n}"),
            ("padding-top",    f"pt-{n}"),
            ("padding-bottom", f"pb-{n}"),
            ("padding-left",   f"pl-{n}"),
            ("padding-right",  f"pr-{n}"),
            ("margin",         f"m-{n}"),
            ("margin-top",     f"mt-{n}"),
            ("margin-bottom",  f"mb-{n}"),
            ("margin-left",    f"ml-{n}"),
            ("margin-right",   f"mr-{n}"),
            ("top",            f"top-{n}"),
            ("bottom",         f"bottom-{n}"),
            ("left",           f"left-{n}"),
            ("right",          f"right-{n}"),
            ("gap",            f"gap-{n}"),
            ("width",          f"w-{n}"),
            ("height",         f"h-{n}"),
        ]:
            rules.append(f".{cls}{{" + prop + f":{v}px!important}}")
        # combined
        rules.append(f".px-{n}{{padding-left:{v}px!important;padding-right:{v}px!important}}")
        rules.append(f".py-{n}{{padding-top:{v}px!important;padding-bottom:{v}px!important}}")
        rules.append(f".mx-{n}{{margin-left:{v}px!important;margin-right:{v}px!important}}")
        rules.append(f".my-{n}{{margin-top:{v}px!important;margin-bottom:{v}px!important}}")

    # ── Fractional widths / heights ──────────────────────────────────────────
    for num, den in [(1,2),(1,3),(2,3),(1,4),(3,4),(1,5),(2,5),(3,5),(4,5),
                     (1,6),(5,6),(1,12),(5,12),(7,12),(11,12)]:
        pct = round(num / den * 100, 4)
        safe = f"{num}\\\\/{den}"
        rules += [f".w-{safe}{{width:{pct}%!important}}",
                  f".h-{safe}{{height:{pct}%!important}}"]

    # ── Named sizes ──────────────────────────────────────────────────────────
    for cls, val in [("full","100%"),("screen","100vw"),("auto","auto"),
                     ("min","min-content"),("max","max-content"),("fit","fit-content")]:
        rules += [f".w-{cls}{{width:{val}!important}}",
                  f".h-{cls}{{height:{val}!important}}"]

    # ── Font sizes ───────────────────────────────────────────────────────────
    for cls, val in [("xs","12px"),("sm","14px"),("base","16px"),("lg","18px"),
                     ("xl","20px"),("2xl","24px"),("3xl","30px"),("4xl","36px"),
                     ("5xl","48px"),("6xl","60px"),("7xl","72px"),("8xl","96px"),("9xl","128px")]:
        rules.append(f".text-{cls}{{font-size:{val}!important}}")

    # ── Font weights ─────────────────────────────────────────────────────────
    for cls, val in [("thin","100"),("extralight","200"),("light","300"),("normal","400"),
                     ("medium","500"),("semibold","600"),("bold","700"),("extrabold","800"),("black","900")]:
        rules.append(f".font-{cls}{{font-weight:{val}!important}}")

    # ── Tailwind colour palette ───────────────────────────────────────────────
    PALETTE: dict[str, str] = {
        "white": "255,255,255", "black": "0,0,0",
        "slate-50":"248,250,252","slate-100":"241,245,249","slate-200":"226,232,240",
        "slate-300":"203,213,225","slate-400":"148,163,184","slate-500":"100,116,139",
        "slate-600":"71,85,105","slate-700":"51,65,85","slate-800":"30,41,59","slate-900":"15,23,42",
        "gray-50":"249,250,251","gray-100":"243,244,246","gray-200":"229,231,235",
        "gray-300":"209,213,219","gray-400":"156,163,175","gray-500":"107,114,128",
        "gray-600":"75,85,99","gray-700":"55,65,81","gray-800":"31,41,55","gray-900":"17,24,39",
        "zinc-800":"39,39,42","zinc-900":"24,24,27","neutral-800":"38,38,38","neutral-900":"23,23,23",
        "stone-800":"41,37,36","stone-900":"28,25,23",
        "red-400":"248,113,113","red-500":"239,68,68","red-600":"220,38,38","red-700":"185,28,28",
        "orange-400":"251,146,60","orange-500":"249,115,22","orange-600":"234,88,12",
        "amber-400":"251,191,36","amber-500":"245,158,11","amber-600":"217,119,6",
        "yellow-400":"250,204,21","yellow-500":"234,179,8",
        "lime-400":"163,230,53","lime-500":"132,204,22",
        "green-400":"74,222,128","green-500":"34,197,94","green-600":"22,163,74","green-700":"21,128,61",
        "emerald-500":"16,185,129","emerald-600":"5,150,105",
        "teal-400":"45,212,191","teal-500":"20,184,166","teal-600":"13,148,136",
        "cyan-400":"34,211,238","cyan-500":"6,182,212","cyan-600":"8,145,178",
        "sky-400":"56,189,248","sky-500":"14,165,233","sky-600":"2,132,199",
        "blue-400":"96,165,250","blue-500":"59,130,246","blue-600":"37,99,235","blue-700":"29,78,216",
        "indigo-400":"129,140,248","indigo-500":"99,102,241","indigo-600":"79,70,229",
        "violet-400":"167,139,250","violet-500":"139,92,246","violet-600":"124,58,237",
        "purple-400":"192,132,252","purple-500":"168,85,247","purple-600":"147,33,218",
        "fuchsia-500":"217,70,239","fuchsia-600":"192,38,211",
        "pink-400":"244,114,182","pink-500":"236,72,153","pink-600":"219,39,119",
        "rose-400":"251,113,133","rose-500":"244,63,94","rose-600":"225,29,72",
    }
    for name, rgb in PALETTE.items():
        rules += [
            f".bg-{name}{{background-color:rgb({rgb})!important}}",
            f".text-{name}{{color:rgb({rgb})!important}}",
            f".border-{name}{{border-color:rgb({rgb})!important}}",
            f".ring-{name}{{--tw-ring-color:rgb({rgb})!important}}",
            f".from-{name}{{--tw-gradient-from:rgb({rgb})!important}}",
            f".to-{name}{{--tw-gradient-to:rgb({rgb})!important}}",
            f".via-{name}{{--tw-gradient-stops:var(--tw-gradient-from),rgb({rgb}),var(--tw-gradient-to)!important}}",
        ]

    # ── Flex / Grid / Layout ─────────────────────────────────────────────────
    rules += [
        ".flex{display:flex!important}", ".inline-flex{display:inline-flex!important}",
        ".flex-col{flex-direction:column!important}", ".flex-row{flex-direction:row!important}",
        ".flex-col-reverse{flex-direction:column-reverse!important}",
        ".flex-wrap{flex-wrap:wrap!important}", ".flex-nowrap{flex-wrap:nowrap!important}",
        ".items-start{align-items:flex-start!important}",
        ".items-center{align-items:center!important}",
        ".items-end{align-items:flex-end!important}",
        ".items-stretch{align-items:stretch!important}",
        ".items-baseline{align-items:baseline!important}",
        ".justify-start{justify-content:flex-start!important}",
        ".justify-center{justify-content:center!important}",
        ".justify-end{justify-content:flex-end!important}",
        ".justify-between{justify-content:space-between!important}",
        ".justify-around{justify-content:space-around!important}",
        ".justify-evenly{justify-content:space-evenly!important}",
        ".self-auto{align-self:auto!important}", ".self-start{align-self:flex-start!important}",
        ".self-center{align-self:center!important}", ".self-end{align-self:flex-end!important}",
        ".self-stretch{align-self:stretch!important}",
        ".flex-1{flex:1 1 0%!important}", ".flex-auto{flex:1 1 auto!important}",
        ".flex-none{flex:none!important}", ".flex-shrink-0{flex-shrink:0!important}",
        ".shrink-0{flex-shrink:0!important}", ".shrink{flex-shrink:1!important}",
        ".grow{flex-grow:1!important}", ".grow-0{flex-grow:0!important}",
        ".grid{display:grid!important}", ".inline-grid{display:inline-grid!important}",
        ".place-items-center{place-items:center!important}",
        ".place-content-center{place-content:center!important}",
        ".col-auto{grid-column:auto!important}", ".row-auto{grid-row:auto!important}",
    ]
    for n in range(1, 13):
        rules += [
            f".grid-cols-{n}{{grid-template-columns:repeat({n},minmax(0,1fr))!important}}",
            f".col-span-{n}{{grid-column:span {n}/span {n}!important}}",
            f".row-span-{n}{{grid-row:span {n}/span {n}!important}}",
            f".col-start-{n}{{grid-column-start:{n}!important}}",
            f".row-start-{n}{{grid-row-start:{n}!important}}",
            f".gap-x-{n*4 if n<7 else n*4}{{column-gap:{n*4}px!important}}",
            f".gap-y-{n*4 if n<7 else n*4}{{row-gap:{n*4}px!important}}",
        ]

    # ── Position ─────────────────────────────────────────────────────────────
    rules += [
        ".static{position:static!important}", ".relative{position:relative!important}",
        ".absolute{position:absolute!important}", ".fixed{position:fixed!important}",
        ".sticky{position:sticky!important}", ".inset-0{inset:0!important}",
        ".inset-auto{inset:auto!important}", ".inset-x-0{left:0!important;right:0!important}",
        ".inset-y-0{top:0!important;bottom:0!important}",
    ]

    # ── Display ──────────────────────────────────────────────────────────────
    rules += [
        ".block{display:block!important}", ".inline-block{display:inline-block!important}",
        ".inline{display:inline!important}", ".hidden{display:none!important}",
        ".table{display:table!important}", ".table-cell{display:table-cell!important}",
        ".overflow-hidden{overflow:hidden!important}", ".overflow-auto{overflow:auto!important}",
        ".overflow-scroll{overflow:scroll!important}", ".overflow-visible{overflow:visible!important}",
        ".overflow-x-hidden{overflow-x:hidden!important}", ".overflow-y-hidden{overflow-y:hidden!important}",
    ]

    # ── Typography ───────────────────────────────────────────────────────────
    rules += [
        ".text-left{text-align:left!important}", ".text-center{text-align:center!important}",
        ".text-right{text-align:right!important}", ".text-justify{text-align:justify!important}",
        ".uppercase{text-transform:uppercase!important}", ".lowercase{text-transform:lowercase!important}",
        ".capitalize{text-transform:capitalize!important}", ".normal-case{text-transform:none!important}",
        ".italic{font-style:italic!important}", ".not-italic{font-style:normal!important}",
        ".underline{text-decoration-line:underline!important}",
        ".line-through{text-decoration-line:line-through!important}",
        ".no-underline{text-decoration-line:none!important}",
        ".leading-none{line-height:1!important}", ".leading-tight{line-height:1.25!important}",
        ".leading-snug{line-height:1.375!important}", ".leading-normal{line-height:1.5!important}",
        ".leading-relaxed{line-height:1.625!important}", ".leading-loose{line-height:2!important}",
        ".tracking-tighter{letter-spacing:-0.05em!important}",
        ".tracking-tight{letter-spacing:-0.025em!important}",
        ".tracking-normal{letter-spacing:0!important}",
        ".tracking-wide{letter-spacing:0.025em!important}",
        ".tracking-wider{letter-spacing:0.05em!important}",
        ".tracking-widest{letter-spacing:0.1em!important}",
        ".truncate{overflow:hidden!important;text-overflow:ellipsis!important;white-space:nowrap!important}",
        ".whitespace-nowrap{white-space:nowrap!important}",
        ".whitespace-pre{white-space:pre!important}", ".whitespace-normal{white-space:normal!important}",
        ".break-words{overflow-wrap:break-word!important}", ".break-all{word-break:break-all!important}",
        ".antialiased{-webkit-font-smoothing:antialiased!important}",
        ".subpixel-antialiased{-webkit-font-smoothing:auto!important}",
    ]

    # ── Max-width presets ─────────────────────────────────────────────────────
    for cls, val in [("xs","320px"),("sm","384px"),("md","448px"),("lg","512px"),
                     ("xl","576px"),("2xl","672px"),("3xl","768px"),("4xl","896px"),
                     ("5xl","1024px"),("6xl","1152px"),("7xl","1280px"),
                     ("prose","65ch"),("full","100%"),("none","none"),
                     ("screen-sm","640px"),("screen-md","768px"),
                     ("screen-lg","1024px"),("screen-xl","1280px"),("screen-2xl","1536px")]:
        rules.append(f".max-w-{cls}{{max-width:{val}!important}}")

    # ── Border radius ─────────────────────────────────────────────────────────
    for cls, val in [("none","0"),("sm","2px"),("","4px"),("md","6px"),("lg","8px"),
                     ("xl","12px"),("2xl","16px"),("3xl","24px"),("full","9999px")]:
        suffix = f"-{cls}" if cls else ""
        rules += [
            f".rounded{suffix}{{border-radius:{val}!important}}",
            f".rounded-t{suffix}{{border-top-left-radius:{val}!important;border-top-right-radius:{val}!important}}",
            f".rounded-b{suffix}{{border-bottom-left-radius:{val}!important;border-bottom-right-radius:{val}!important}}",
            f".rounded-l{suffix}{{border-top-left-radius:{val}!important;border-bottom-left-radius:{val}!important}}",
            f".rounded-r{suffix}{{border-top-right-radius:{val}!important;border-bottom-right-radius:{val}!important}}",
        ]

    # ── Border ────────────────────────────────────────────────────────────────
    for w in [0, 1, 2, 4, 8]:
        suffix = f"-{w}" if w != 1 else ""
        if w == 1:
            rules += [
                ".border{border-width:1px!important;border-style:solid!important}",
                ".border-t{border-top-width:1px!important;border-top-style:solid!important}",
                ".border-b{border-bottom-width:1px!important;border-bottom-style:solid!important}",
                ".border-l{border-left-width:1px!important;border-left-style:solid!important}",
                ".border-r{border-right-width:1px!important;border-right-style:solid!important}",
            ]
        else:
            rules.append(f".border-{w}{{border-width:{w}px!important;border-style:solid!important}}")

    # ── Opacity / z-index ─────────────────────────────────────────────────────
    for z in [0, 10, 20, 30, 40, 50, 100]:
        rules.append(f".z-{z}{{z-index:{z}!important}}")
    for o in range(0, 101, 5):
        rules.append(f".opacity-{o}{{opacity:{o/100:.2f}!important}}")

    # ── Shadows ───────────────────────────────────────────────────────────────
    rules += [
        ".shadow-sm{box-shadow:0 1px 2px 0 rgb(0 0 0/0.05)!important}",
        ".shadow{box-shadow:0 1px 3px 0 rgb(0 0 0/0.1),0 1px 2px -1px rgb(0 0 0/0.1)!important}",
        ".shadow-md{box-shadow:0 4px 6px -1px rgb(0 0 0/0.1),0 2px 4px -2px rgb(0 0 0/0.1)!important}",
        ".shadow-lg{box-shadow:0 10px 15px -3px rgb(0 0 0/0.1),0 4px 6px -4px rgb(0 0 0/0.1)!important}",
        ".shadow-xl{box-shadow:0 20px 25px -5px rgb(0 0 0/0.1),0 8px 10px -6px rgb(0 0 0/0.1)!important}",
        ".shadow-2xl{box-shadow:0 25px 50px -12px rgb(0 0 0/0.25)!important}",
        ".shadow-inner{box-shadow:inset 0 2px 4px 0 rgb(0 0 0/0.05)!important}",
        ".shadow-none{box-shadow:none!important}",
    ]

    # ── Miscellaneous ─────────────────────────────────────────────────────────
    rules += [
        ".object-cover{object-fit:cover!important}", ".object-contain{object-fit:contain!important}",
        ".object-fill{object-fit:fill!important}", ".object-center{object-position:center!important}",
        ".pointer-events-none{pointer-events:none!important}",
        ".select-none{user-select:none!important}",
        ".cursor-pointer{cursor:pointer!important}",
        ".resize-none{resize:none!important}",
        ".appearance-none{appearance:none!important}",
        ".box-border{box-sizing:border-box!important}",
        ".box-content{box-sizing:content-box!important}",
        ".list-none{list-style-type:none!important}",
        ".list-disc{list-style-type:disc!important}",
        ".list-decimal{list-style-type:decimal!important}",
        ".mx-auto{margin-left:auto!important;margin-right:auto!important}",
        ".my-auto{margin-top:auto!important;margin-bottom:auto!important}",
        ".mt-auto{margin-top:auto!important}", ".mb-auto{margin-bottom:auto!important}",
        ".ml-auto{margin-left:auto!important}", ".mr-auto{margin-right:auto!important}",
        ".bg-transparent{background-color:transparent!important}",
        ".bg-gradient-to-r{background-image:linear-gradient(to right,var(--tw-gradient-stops))!important}",
        ".bg-gradient-to-l{background-image:linear-gradient(to left,var(--tw-gradient-stops))!important}",
        ".bg-gradient-to-b{background-image:linear-gradient(to bottom,var(--tw-gradient-stops))!important}",
        ".bg-gradient-to-t{background-image:linear-gradient(to top,var(--tw-gradient-stops))!important}",
        ".bg-gradient-to-br{background-image:linear-gradient(to bottom right,var(--tw-gradient-stops))!important}",
        ".bg-gradient-to-bl{background-image:linear-gradient(to bottom left,var(--tw-gradient-stops))!important}",
        ".bg-gradient-to-tr{background-image:linear-gradient(to top right,var(--tw-gradient-stops))!important}",
        ".bg-gradient-to-tl{background-image:linear-gradient(to top left,var(--tw-gradient-stops))!important}",
        ":root{--tw-gradient-from:#fff;--tw-gradient-to:rgb(255 255 255/0);--tw-gradient-stops:var(--tw-gradient-from),var(--tw-gradient-to)}",
    ]

    return "\n".join(rules)


# Pre-build the static CSS (done once at import time)
_STATIC_TAILWIND_CSS = _build_tailwind_override_css()


def _inject_css(html: str) -> str:
    """
    Inject:
      1. The pre-built static Tailwind utility CSS.
      2. Dynamic arbitrary-value rules extracted from this specific HTML.
    """
    ARBITRARY_PROPS = {
        'text': 'font-size', 'w': 'width', 'h': 'height',
        'max-w': 'max-width', 'min-w': 'min-width',
        'max-h': 'max-height', 'min-h': 'min-height',
        'p': 'padding', 'px': 'padding-inline', 'py': 'padding-block',
        'pt': 'padding-top', 'pb': 'padding-bottom',
        'pl': 'padding-left', 'pr': 'padding-right',
        'm': 'margin', 'mt': 'margin-top', 'mb': 'margin-bottom',
        'ml': 'margin-left', 'mr': 'margin-right',
        'gap': 'gap', 'gap-x': 'column-gap', 'gap-y': 'row-gap',
        'leading': 'line-height', 'tracking': 'letter-spacing',
        'rounded': 'border-radius', 'border': 'border-width',
        'top': 'top', 'bottom': 'bottom', 'left': 'left', 'right': 'right',
        'z': 'z-index', 'opacity': 'opacity', 'translate-x': 'transform',
        'basis': 'flex-basis', 'col-span': 'grid-column',
    }
    dynamic_rules: list[str] = []
    seen: set[str] = set()
    for prefix, value in re.findall(r'([\w-]+)-\[([^\]]+)\]', html):
        key = f"{prefix}-[{value}]"
        if key in seen or prefix not in ARBITRARY_PROPS:
            continue
        seen.add(key)
        # Escape special chars for CSS selector
        cls = re.sub(r'([\[\]%#().,])', lambda m: '\\' + m.group(), key)
        prop = ARBITRARY_PROPS[prefix]
        dynamic_rules.append(f".{cls}{{{prop}:{value}!important}}")

    full_css = (
        "<style>\n"
        + _STATIC_TAILWIND_CSS + "\n"
        + "\n".join(dynamic_rules) + "\n"
        + "</style>\n"
    )
    if "</head>" in html:
        return html.replace("</head>", full_css + "</head>", 1)
    return full_css + html


# JS text extractor: uses WeakSet absorption to prevent double text rendering.
_JS_EXTRACT_TEXT = """
() => {
    function rgbParse(s) {
        if (!s) return null;
        const m = s.match(/rgba?\\((\\d+),\\s*(\\d+),\\s*(\\d+)(?:,\\s*([\\d.]+))?\\)/);
        if (!m) return null;
        const a = m[4] !== undefined ? parseFloat(m[4]) : 1;
        if (a < 0.05) return null;
        return [+m[1], +m[2], +m[3]];
    }
    function firstFont(s) {
        return (s || '').split(',')[0].replace(/['"]/g, '').trim() || 'Arial';
    }
    function isIconFont(cs) {
        const f = (cs.fontFamily || '').toLowerCase();
        return f.includes('material') || f.includes('fontawesome') ||
               f.includes('font awesome') || f.includes('icons') ||
               f.includes('feather');
    }

    let slideEl = (
        document.querySelector('.slide') || document.querySelector('.ppt-slide') ||
        (() => {
            for (const d of document.querySelectorAll('div')) {
                const r = d.getBoundingClientRect();
                if (r.width >= 600 && r.height >= 300) return d;
            }
            return null;
        })() || document.body
    );
    const sRect = slideEl.getBoundingClientRect();
    const offX = sRect.left, offY = sRect.top;
    const SW = sRect.width || 1280;
    const SH = sRect.height || 720;

    const INLINE = new Set(['span','a','strong','em','b','i','u','s','mark',
                             'code','kbd','label','sup','sub','cite','abbr']);
    const absorbed = new WeakSet();
    const results = [];

    for (const el of slideEl.querySelectorAll('*')) {
        if (absorbed.has(el)) continue;

        const rect = el.getBoundingClientRect();
        if (rect.width < 1 || rect.height < 1) continue;
        if (rect.right <= offX || rect.bottom <= offY || rect.left >= offX + SW || rect.top >= offY + SH) continue;

        const cs = window.getComputedStyle(el);
        if (isIconFont(cs)) { absorbed.add(el); continue; }
        if (cs.visibility === 'hidden' || cs.display === 'none' || cs.opacity === '0') continue;

        const runs = [];
        for (const child of el.childNodes) {
            if (child.nodeType === 3) {
                const t = child.textContent.replace(/\\s+/g, ' ').trim();
                if (!t) continue;
                runs.push({
                    text: t,
                    color: rgbParse(cs.color),
                    fontSize: parseFloat(cs.fontSize),
                    fontWeight: parseInt(cs.fontWeight) || 400,
                    fontFamily: firstFont(cs.fontFamily),
                    italic: cs.fontStyle === 'italic',
                    underline: cs.textDecoration.includes('underline'),
                });
            } else if (child.nodeType === 1 && INLINE.has(child.tagName.toLowerCase())) {
                absorbed.add(child);
                for (const desc of child.querySelectorAll('*')) absorbed.add(desc);

                const ccs = window.getComputedStyle(child);
                if (isIconFont(ccs)) continue;
                if (ccs.visibility === 'hidden' || ccs.display === 'none' || ccs.opacity === '0') continue;
                const t = (child.innerText || child.textContent || '').replace(/\\s+/g, ' ').trim();
                if (!t) continue;
                runs.push({
                    text: t,
                    color: rgbParse(ccs.color) || rgbParse(cs.color),
                    fontSize: parseFloat(ccs.fontSize),
                    fontWeight: parseInt(ccs.fontWeight) || 400,
                    fontFamily: firstFont(ccs.fontFamily),
                    italic: ccs.fontStyle === 'italic',
                    underline: ccs.textDecoration.includes('underline'),
                });
            }
        }
        if (!runs.length) continue;

        results.push({
            tag: el.tagName.toLowerCase(),
            x: Math.round(rect.left - offX),
            y: Math.round(rect.top - offY),
            w: Math.round(rect.width),
            h: Math.round(rect.height),
            runs,
            textAlign: cs.textAlign === 'start' ? 'left' : (cs.textAlign || 'left'),
        });
    }
    return { elements: results, slideW: Math.round(SW), slideH: Math.round(SH) };
}
"""

def _align(text_align: str) -> PP_ALIGN:
    return {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.DISTRIBUTE}.get(text_align, PP_ALIGN.LEFT)

def _add_textbox(slide, x, y, w, h, runs: list[dict], text_align: str,
                 default_color: Optional[RGBColor] = None):
    if not runs or not any(r.get('text') for r in runs):
        return
    w = max(w, 20); h = max(h, 10)
    txbox = slide.shapes.add_textbox(_px(x), _px(y), _px(w), _px(h))
    tf = txbox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = _align(text_align)
    for i, rd in enumerate(runs):
        text = rd.get('text', '').strip()
        if not text: continue
        if i > 0 and not text.startswith(' '):
            prev = runs[i - 1].get('text', '')
            if prev and not prev.endswith(' '): text = ' ' + text
        run = p.add_run(); run.text = text
        run.font.size = Pt(rd.get('fontSize', 16) * 0.75)
        run.font.bold = rd.get('fontWeight', 400) >= 600
        run.font.italic = rd.get('italic', False)
        run.font.underline = rd.get('underline', False)
        color = _rgba_to_rgb(rd.get('color'))
        run.font.color.rgb = color if color else (default_color or RGBColor(0,0,0))
        FONT_MAP = {
            'Source Code Pro': 'Courier New', 'Courier': 'Courier New',
            'Times New Roman': 'Times New Roman', 'Georgia': 'Georgia',
            'Roboto': 'Calibri', 'Roboto Flex': 'Calibri', 'Inter': 'Calibri',
            'Arial': 'Arial', 'Helvetica': 'Arial',
        }
        run.font.name = FONT_MAP.get(rd.get('fontFamily'), 'Calibri')

def _add_image_shape(slide, x, y, w, h, src: str):
    try:
        if src.startswith('data:'):
            import base64; header, data = src.split(',', 1)
            img_data = BytesIO(base64.b64decode(data))
        elif src.startswith('http'):
            img_data = BytesIO(requests.get(src, timeout=8).content)
        elif os.path.exists(src):
            img_data = BytesIO(open(src, 'rb').read())
        else: return
        slide.shapes.add_picture(img_data, _px(x), _px(y), _px(w), _px(h))
    except Exception as e: print(f"[Image] error: {e}")

def _screenshot_to_image(screenshot_bytes: bytes) -> BytesIO:
    img = PILImage.open(BytesIO(screenshot_bytes))
    if img.size != (BROWSER_W, BROWSER_H):
        img = img.resize((BROWSER_W, BROWSER_H), PILImage.LANCZOS)
    buf = BytesIO(); img.save(buf, format='PNG'); buf.seek(0)
    return buf

def _split_slides(html: str) -> list[str]:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(html, 'html.parser')
    head = str(soup.find('head')) if soup.find('head') else "<head></head>"
    slides = soup.find_all('div', class_=lambda c: c and ('slide' in c.split() or 'ppt-slide' in c.split()))
    if not slides: return [html]
    return [f"<!DOCTYPE html><html>{head}<body>{s}</body></html>" for s in slides]

class PresentationCreator:
    def __init__(self, font_name: str = 'Calibri', timeout_ms: int = 4000, overlay_text: bool = True):
        self.font_name = font_name; self.timeout_ms = timeout_ms; self.overlay_text = overlay_text

    def generate_filename(self, filename: str = None) -> str:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = filename or f"presentation_{timestamp}"
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        return filename

    async def create_presentation(self, content: str, filename: str = None) -> BytesIO:
        slide_docs = _split_slides(content); prs = Presentation()
        prs.slide_width, prs.slide_height = PPTX_W, PPTX_H
        async with async_playwright() as pw:
            browser = await pw.chromium.launch(args=['--no-sandbox', '--disable-dev-shm-usage'])
            for doc in slide_docs:
                await self._render_slide(browser, prs, doc)
            await browser.close()
        buf = BytesIO(); prs.save(buf); buf.seek(0); return buf

    async def _render_slide(self, browser, prs: Presentation, html: str):
        patched = _inject_css(html)
        tf = tempfile.NamedTemporaryFile(suffix='.html', delete=False, mode='w', encoding='utf-8')
        tf.write(patched); tf.close()
        try:
            page = await browser.new_page(viewport={'width': BROWSER_W, 'height': BROWSER_H})
            await page.goto(f'file://{tf.name}', wait_until='networkidle')
            await page.wait_for_timeout(self.timeout_ms)
            text_data = await page.evaluate(_JS_EXTRACT_TEXT) if self.overlay_text else None
            await page.evaluate("""() => {
                const isIcon = (el) => {
                    const f = (window.getComputedStyle(el).fontFamily || '').toLowerCase();
                    return f.includes('material') || f.includes('fontawesome') || f.includes('font awesome') || f.includes('icons');
                };
                document.querySelectorAll('*').forEach(el => { if (isIcon(el)) el.dataset.pptxKeep = 'true'; });
                const s = document.createElement('style');
                s.textContent = `*:not([data-pptx-keep]):not([data-pptx-keep] *) {
                    color: transparent !important; -webkit-text-fill-color: transparent !important;
                    fill: transparent !important; text-shadow: none !important;
                }`;
                document.head.appendChild(s);
            }""")
            await page.wait_for_timeout(200)
            screenshot = await page.screenshot(type='png', clip={'x':0,'y':0,'width':BROWSER_W,'height':BROWSER_H})
            await page.close()
        finally: os.unlink(tf.name)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(_screenshot_to_image(screenshot), 0, 0, PPTX_W, PPTX_H)
        if text_data:
            for el in text_data.get('elements', []):
                _add_textbox(slide, el['x'], el['y'], el['w'], el['h'], el['runs'], el['textAlign'], RGBColor(0,0,0))
